from flask import Blueprint, render_template, send_file, jsonify
from models.file import File
from models.storage import Storage
from database import db_session
import os
import json
import boto3
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import io
import mimetypes

preview_bp = Blueprint('preview', __name__)

def get_storage_client(storage):
    """根据存储池类型获取对应的存储客户端"""
    config = json.loads(storage.config)
    if storage.type == 'local':
        return LocalStorage(config['path'])
    elif storage.type == 's3':
        return S3Storage(
            access_key=config['access_key'],
            secret_key=config['secret_key'],
            bucket=config['bucket'],
            region=config['region']
        )
    raise ValueError(f'不支持的存储类型: {storage.type}')

class LocalStorage:
    def __init__(self, base_path):
        self.base_path = base_path
    
    def get_file_path(self, filename):
        return os.path.join(self.base_path, filename)
    
    def get_file_content(self, filename):
        file_path = self.get_file_path(filename)
        with open(file_path, 'rb') as f:
            return f.read()

class S3Storage:
    def __init__(self, access_key, secret_key, bucket, region):
        self.s3 = boto3.client(
            's3',
            aws_access_key_id=access_key,
            aws_secret_access_key=secret_key,
            region_name=region
        )
        self.bucket = bucket
    
    def get_file_content(self, filename):
        response = self.s3.get_object(Bucket=self.bucket, Key=filename)
        return response['Body'].read()

@preview_bp.route('/preview/<path:filename>')
def preview_file(filename):
    """预览文件"""
    # 获取当前活动的存储池
    storage = db_session.query(Storage).filter_by(is_active=1).first()
    if not storage:
        return render_template('error.html', message='没有活动的存储池'), 404
    
    # 查找文件记录
    file_record = db_session.query(File).filter_by(
        storage_id=storage.id,
        name=filename
    ).first()
    
    if not file_record:
        return render_template('error.html', message='文件不存在'), 404
    
    try:
        # 获取存储客户端
        storage_client = get_storage_client(storage)
        
        # 获取文件内容
        file_content = storage_client.get_file_content(file_record.path)
        
        # 根据文件类型处理预览
        file_type = file_record.type.lower()
        
        if 'word' in file_type or filename.endswith('.docx'):
            return preview_word(file_content, filename)
        elif 'excel' in file_type or filename.endswith(('.xlsx', '.xls')):
            return preview_excel(file_content, filename)
        elif 'powerpoint' in file_type or filename.endswith(('.pptx', '.ppt')):
            return preview_ppt(file_content, filename)
        elif 'image' in file_type:
            return preview_image(file_content, filename)
        else:
            return render_template('error.html', message='不支持预览此类型的文件'), 400
    
    except Exception as e:
        return render_template('error.html', message=str(e)), 500

def preview_word(file_content, filename):
    """预览 Word 文档"""
    doc = Document(io.BytesIO(file_content))
    content = []
    
    # 提取文本内容
    for para in doc.paragraphs:
        if para.text.strip():
            content.append({
                'type': 'paragraph',
                'text': para.text
            })
    
    # 提取表格内容
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [cell.text for cell in row.cells]
            table_data.append(row_data)
        content.append({
            'type': 'table',
            'data': table_data
        })
    
    return render_template('preview/word.html', 
                         filename=filename,
                         content=content)

def preview_excel(file_content, filename):
    """预览 Excel 文档"""
    wb = load_workbook(io.BytesIO(file_content))
    sheets = []
    
    # 提取每个工作表的内容
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        data = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell for cell in row):  # 只添加非空行
                data.append(list(row))
        sheets.append({
            'name': sheet_name,
            'data': data
        })
    
    return render_template('preview/excel.html',
                         filename=filename,
                         sheets=sheets)

def preview_ppt(file_content, filename):
    """预览 PowerPoint 文档"""
    prs = Presentation(io.BytesIO(file_content))
    slides = []
    
    # 提取每张幻灯片的内容
    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)
        slides.append(slide_content)
    
    return render_template('preview/ppt.html',
                         filename=filename,
                         slides=slides)

def preview_image(file_content, filename):
    """预览图片"""
    # 将图片转换为 base64 编码
    import base64
    image = Image.open(io.BytesIO(file_content))
    buffered = io.BytesIO()
    image.save(buffered, format=image.format)
    img_str = base64.b64encode(buffered.getvalue()).decode()
    
    return render_template('preview/image.html',
                         filename=filename,
                         image_data=img_str) 